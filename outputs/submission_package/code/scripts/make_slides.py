#!/usr/bin/env python
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("outputs/presentation_stop_verify_explore.pptx")

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)
FONT = "Arial"
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(90, 90, 90)
LIGHT_GRAY = RGBColor(235, 235, 235)
BLUE = RGBColor(37, 99, 235)


def main():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H

    add_title_slide(prs)
    add_motivation(prs)
    add_problem_setting(prs)
    add_method(prs)
    add_oracle(prs)
    add_controllers(prs)
    add_main_result(prs)
    add_model_comparison(prs)
    add_limitations(prs)
    add_conclusion(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


def add_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(12.0), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = BLACK


def add_footer(slide, n):
    box = slide.shapes.add_textbox(Inches(6.45), Inches(7.15), Inches(0.5), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.color.rgb = BLACK


def add_bullets(slide, bullets, x=0.85, y=1.75, w=11.7, font_size=27, gap=0.14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.2))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(16 + gap * 10)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12), Inches(2.1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Stop, Verify, or Explore?"
    r.font.name = FONT
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = BLACK

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(3.15), Inches(11.8), Inches(1.1))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "State-Aware Value-of-Compute Routing for Autoregressive LLM Reasoning"
    r.font.name = FONT
    r.font.size = Pt(25)
    r.font.color.rgb = BLACK

    meta = slide.shapes.add_textbox(Inches(0.78), Inches(5.25), Inches(10.5), Inches(0.8))
    p = meta.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "GSM8K · Qwen2.5-1.5B / 3B · STOP / VERIFY / SC-3"
    r.font.name = FONT
    r.font.size = Pt(22)
    r.font.color.rgb = GRAY
    add_footer(slide, 1)


def add_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Motivation")
    add_bullets(
        slide,
        [
            "Reasoning 문제마다 필요한 test-time compute가 다르다.",
            "항상 더 많이 생성하면 accuracy는 오를 수 있지만 token cost가 커진다.",
            "핵심 질문: 현재 reasoning state를 보고 extra compute를 쓸지 결정할 수 있는가?",
            "목표: accuracy–compute trade-off 개선",
        ],
        font_size=29,
    )
    add_footer(slide, 2)


def add_problem_setting(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Problem Setting")
    add_flow(slide)
    add_bullets(
        slide,
        [
            "Initial answer 생성 후 한 번만 routing하는 one-step setting",
            "Multi-step RL / MDP가 아니라 post-hoc routing 문제",
            "Controller는 STOP, VERIFY, SC-3 중 정확히 하나를 선택",
        ],
        x=0.9,
        y=4.55,
        font_size=25,
    )
    add_footer(slide, 3)


def add_flow(slide):
    labels = ["Question", "Initial\nReasoning", "Router", "Final\nAnswer"]
    xs = [0.85, 3.55, 6.25, 9.0]
    y = 2.05
    for i, (label, x) in enumerate(zip(labels, xs)):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.0), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY if i != 2 else RGBColor(220, 235, 255)
        shape.line.color.rgb = BLACK
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.size = Pt(22)
        r.font.bold = i == 2
        if i < len(labels) - 1:
            arr = slide.shapes.add_textbox(Inches(x + 2.08), Inches(y + 0.25), Inches(0.7), Inches(0.4))
            pr = arr.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.CENTER
            rr = pr.add_run()
            rr.text = "→"
            rr.font.name = FONT
            rr.font.size = Pt(34)


def add_method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Method: STOP / VERIFY / SC-3")
    rows = [
        ["STOP", "initial answer 그대로 사용", "0"],
        ["VERIFY", "문제를 다시 풀고 기존 답과 비교", "1 generation"],
        ["SC-3", "추가 독립 풀이 2개 + plurality vote", "2 generations"],
    ]
    add_table(slide, ["Action", "Decision", "Extra compute"], rows, x=0.8, y=1.75, w=11.7, h=2.5, font_size=20)
    add_bullets(
        slide,
        [
            "GSM8K 답변은 `Final answer: <answer>` 형식으로 추출",
            "SC-3 tie-break는 initial answer",
        ],
        x=0.95,
        y=4.65,
        font_size=25,
    )
    add_footer(slide, 4)


def add_oracle(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Oracle Router")
    formula = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.4), Inches(0.8))
    p = formula.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "utility(action) = correct(action) − λ · normalized_cost(action)"
    r.font.name = FONT
    r.font.size = Pt(29)
    r.font.bold = True
    r.font.color.rgb = BLACK

    add_bullets(
        slide,
        [
            "각 example에서 STOP / VERIFY / SC-3의 실제 rollout 결과를 모두 사용",
            "Oracle action = utility가 가장 큰 action",
            "Learned controller가 도달할 수 있는 upper bound 역할",
        ],
        x=0.95,
        y=3.0,
        font_size=26,
    )
    add_footer(slide, 5)


def add_controllers(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Learned Controllers")
    add_bullets(
        slide,
        [
            "Input/state-aware features: question length, numeric tokens, initial trace length, token counts",
            "Oracle-action classifier: LogisticRegression / RandomForest",
            "Validation calibration: extra-compute confidence threshold 조정",
            "Value predictor: action별 correctness와 cost를 예측 후 utility 최대 action 선택",
        ],
        font_size=26,
    )
    add_footer(slide, 6)


def add_main_result(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Main Result: 1.5B, 300 Examples")
    rows = [
        ["Always STOP", "0.253", "299.68", "0.222"],
        ["Always VERIFY", "0.370", "916.16", "0.275"],
        ["Always SC-3", "0.287", "966.53", "0.187"],
        ["Oracle Router", "0.473", "435.65", "0.428"],
        ["RF calibrated", "0.430", "484.83", "0.380"],
        ["Value predictor", "0.400", "778.28", "0.319"],
    ]
    add_table(slide, ["Method", "Acc.", "Avg tokens", "Utility"], rows, x=0.75, y=1.55, w=11.8, h=3.9, font_size=17)
    note = slide.shapes.add_textbox(Inches(0.95), Inches(6.0), Inches(11.2), Inches(0.55))
    p = note.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Oracle improves STOP accuracy 0.253 → 0.473; learned RF recovers much of the gain."
    r.font.name = FONT
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = BLACK
    add_footer(slide, 7)


def add_model_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Model Comparison: 1.5B vs 3B")
    rows = [
        ["Qwen2.5-1.5B", "0.290", "0.440", "0.320", "0.530", "0.470"],
        ["Qwen2.5-3B", "0.320", "0.540", "0.340", "0.600", "0.540"],
    ]
    add_table(slide, ["Model", "STOP", "VERIFY", "SC-3", "Oracle", "Best Ctrl"], rows, x=0.65, y=1.65, w=12.0, h=2.0, font_size=17)
    add_bullets(
        slide,
        [
            "동일한 100 shuffled GSM8K examples 기준 비교",
            "3B는 STOP / VERIFY / Oracle / best controller 모두 향상",
            "VERIFY gain이 가장 큼: 0.440 → 0.540",
        ],
        x=0.95,
        y=4.3,
        font_size=26,
    )
    add_footer(slide, 8)


def add_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Limitations")
    add_bullets(
        slide,
        [
            "Dataset은 GSM8K 중심: MATH500 등 cross-dataset 검증은 남음",
            "3B는 100 examples라 test split variance가 큼",
            "Entropy / token log-probability feature는 아직 미구현",
            "Oracle labels are imbalanced: 대부분 STOP이 optimal",
            "SC-3만 실험: SC-5 이상에서는 다른 양상이 가능",
        ],
        font_size=26,
    )
    add_footer(slide, 9)


def add_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusion")
    add_bullets(
        slide,
        [
            "Extra compute는 모든 문제에 필요한 것은 아니지만, 일부 문제에서는 큰 gain을 만든다.",
            "Oracle routing은 STOP 대비 큰 accuracy/utility 개선을 보인다.",
            "Learned controller는 calibration과 model choice에 민감하지만 일부 이득을 회복한다.",
            "Multiple models에서도 value-of-compute signal이 유지된다.",
        ],
        font_size=28,
    )
    call = slide.shapes.add_textbox(Inches(0.95), Inches(6.05), Inches(11.0), Inches(0.6))
    p = call.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Takeaway: Stop when enough, verify when useful, explore only selectively."
    r.font.name = FONT
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = BLUE
    add_footer(slide, 10)


def add_table(slide, headers, rows, x, y, w, h, font_size=18):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY
        style_cell(cell, font_size, bold=True)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            style_cell(cell, font_size, bold=False)
    # Make first column wider for method names.
    if len(headers) >= 4:
        table.columns[0].width = Inches(w * 0.36)
        rest = (w * 0.64) / (len(headers) - 1)
        for col in range(1, len(headers)):
            table.columns[col].width = Inches(rest)


def style_cell(cell, font_size, bold=False):
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = BLACK


if __name__ == "__main__":
    main()

